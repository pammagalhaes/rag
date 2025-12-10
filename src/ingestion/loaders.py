import os
from typing import List
from langchain_core.documents import Document
from pypdf import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image




def load_pdf(path: str) -> List[Document]:
    reader = PdfReader(path)
    pages = []
    for i, p in enumerate(reader.pages):
        text = p.extract_text() or ""
        if text.strip():
            pages.append(Document(page_content=text, metadata={"source": os.path.basename(path), "page": i}))
    return pages




def load_txt(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]




def load_md(path: str) -> List[Document]:
    return load_txt(path)




def load_pptx(path: str) -> List[Document]:
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)
        txt = "\n".join(content).strip()
        if txt:
            docs.append(Document(page_content=txt, metadata={"source": os.path.basename(path), "slide": i}))
    return docs




def load_image(path: str) -> List[Document]:
    text = pytesseract.image_to_string(Image.open(path))
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]




def load_file(path: str) -> List[Document]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return load_pdf(path)
    if ext in [".txt"]:
        return load_txt(path)
    if ext in [".md", ".markdown"]:
        return load_md(path)
    if ext in [".pptx"]:
        return load_pptx(path)
    # if ext in [".png", ".jpg", ".jpeg", ".tiff"]:
    #     return load_image(path)
    raise ValueError(f"Unsupported extension: {ext}")